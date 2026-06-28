#!/usr/bin/env python3
"""
hymn_to_presentation.py — Convert a hymn file to HTML, PDF, PPTX, OTP, pro6, or ewsx slides.

Produces one title slide (with full attribution) followed by content slides
showing up to --lines-per-slide (default 3) melody/lyric pairs each.  The
MusiQwik melody line appears immediately above the corresponding lyric line.
Each verse starts on a new slide.

Usage:
    python3 hymn_to_presentation.py --file <hymn_file>
    python3 hymn_to_presentation.py --file <hymn_file> --format pptx
    python3 hymn_to_presentation.py --file <hymn_file> --format pdf -o out.pdf
    python3 hymn_to_presentation.py --file <hymn_file> --format otp
    python3 hymn_to_presentation.py --file <hymn_file> --format pro6
    python3 hymn_to_presentation.py --file <hymn_file> --format ewsx
    python3 hymn_to_presentation.py --file <hymn_file> --lines-per-slide 2

Tradition filtering (stanza-level [Tags: ...] markers in the hymn file):
    python3 hymn_to_presentation.py --file <hymn_file> --include lutheran
    python3 hymn_to_presentation.py --file <hymn_file> --exclude baptist
    python3 hymn_to_presentation.py --file <hymn_file> -I lutheran -I ecumenical -X anglican

    --include / -I  Keep only stanzas whose tradition tags overlap this set.
    --exclude / -X  Drop any stanza whose tradition tags overlap this set.

    Stanzas without a [Tags: ...] marker inherit the file-level Tags: line.
    Both flags are repeatable; omitting both emits all stanzas unchanged.

Dependencies:
    PDF  — playwright  (pip install playwright)  +  Chromium at /opt/pw-browsers
    PPTX — python-pptx (pip install python-pptx)
    HTML/OTP/pro6/ewsx — stdlib only
"""

import argparse
import base64
import io
import re
import sys
import uuid
import zipfile
from collections import defaultdict
from dataclasses import dataclass, field
from html import escape as _xe
from pathlib import Path


# ---------------------------------------------------------------------------
# Hymn-file parsing
# ---------------------------------------------------------------------------

# Matches [Tags: tag1, tag2] stanza-level markers at the start of a verse chunk.
_STANZA_TAG_RE = re.compile(r'^\s*\[Tags:\s*([^\]]+)\]', re.IGNORECASE)


def _extract_section(text: str, header_re: str, stop_re: str) -> str:
    m = re.search(header_re, text, re.MULTILINE | re.IGNORECASE)
    if not m:
        return ''
    start = m.end()
    stop = re.search(stop_re, text[start:], re.MULTILINE)
    return text[start : start + stop.start()] if stop else text[start:]


def parse_hymn_file(path: Path) -> dict:
    text  = path.read_text(encoding='utf-8')
    lines = text.splitlines()
    title = next((l.strip() for l in lines if l.strip()), '')

    tag_match = re.search(r'^Tags:[ \t]*(.+)$', text, re.MULTILINE)
    file_tags = (
        [t.strip() for t in tag_match.group(1).split(',') if t.strip()]
        if tag_match else []
    )

    return {
        'title':       title,
        'file_tags':   file_tags,
        'abc_text':    _extract_section(text, r'^##\s*ABC[^\n]*\n',      r'^##|^#(?!#)'),
        'musiqwik':    _extract_section(text, r'^##\s*Musiquik[^\n]*\n', r'^##|^#(?!#)').strip(),
        'lyrics_text': _extract_section(text, r'^#\s*Lyrics[^\n]*\n',    r'^#(?!#)').strip(),
        'attribution': _extract_section(text, r'^#\s*Citations[^\n]*\n', r'^#(?!#)').strip(),
    }


# ---------------------------------------------------------------------------
# Melody-line splitting
# ---------------------------------------------------------------------------

def _abc_body_lines(abc_text: str) -> list[str]:
    """Non-header, non-directive lines from an ABC block."""
    in_body, result = False, []
    for line in abc_text.splitlines():
        s = line.strip()
        if not in_body:
            if s and not re.match(r'^[A-Za-z]:', s):
                in_body = True
            else:
                continue
        if s and not s.startswith('%'):
            result.append(s)
    return result


def _barlines_per_line(abc_lines: list[str]) -> list[int]:
    """Count measure-ending barlines for each physical ABC body line."""
    counts = []
    for line in abc_lines:
        cleaned = re.sub(r'\|:', '', line)       # begin-repeat doesn't end a measure
        cleaned = re.sub(r'\(\d+', '', cleaned)  # strip triplet/tuplet markers
        n = len(re.findall(r'\|{1,2}|:\|', cleaned))
        if n:
            counts.append(n)
    return counts


def split_musiqwik(musiqwik: str, barlines_per_line: list[int]) -> list[str]:
    """
    Split the MusiQwik string into one segment per melody line.
    Each segment is prefixed with the original clef+time-sig so it is
    self-contained when pasted into a slide.
    """
    if not musiqwik or not barlines_per_line:
        return [musiqwik] if musiqwik else []

    prefix  = musiqwik[:2]   # '&4', '&3', etc.
    content = musiqwik[2:]

    segments: list[str] = []
    pos = 0
    for needed in barlines_per_line:
        found, i = 0, pos
        while i < len(content) and found < needed:
            if content[i] in '.)':
                found += 1
            i += 1
        segments.append(prefix + content[pos:i])
        pos = i

    # Attach any tail to the last segment (shouldn't happen with correct counts)
    if pos < len(content):
        if segments:
            segments[-1] += content[pos:]
        else:
            segments.append(prefix + content[pos:])

    return segments


# ---------------------------------------------------------------------------
# Lyric parsing
# ---------------------------------------------------------------------------

def parse_verses(lyrics_text: str) -> list[tuple[str, list[str]]]:
    """
    Split continuous lyric text into (verse_text, stanza_tags) pairs.

    Each chunk is checked for a leading [Tags: ...] marker; if found, the
    tags are extracted and the marker stripped from the display text.
    Chunks without a marker return an empty tag list (the caller should fall
    back to the file-level Tags: line for filtering purposes).
    """
    parts  = re.split(r'(?<!\w)(\d+\.)\s+', lyrics_text)
    chunks: list[str] = []
    if parts and parts[0].strip():
        chunks.append(parts[0].strip())
    for i in range(1, len(parts) - 1, 2):
        v = parts[i + 1].strip() if (i + 1) < len(parts) else ''
        if v:
            chunks.append(v)

    verses: list[tuple[str, list[str]]] = []
    for chunk in chunks:
        m = _STANZA_TAG_RE.match(chunk)
        if m:
            tags = [t.strip() for t in m.group(1).split(',') if t.strip()]
            text = chunk[m.end():].strip()
        else:
            tags = []
            text = chunk
        verses.append((text, tags))
    return verses


def filter_verses(
    verses: list[tuple[str, list[str]]],
    file_tags: list[str],
    include: list[str],
    exclude: list[str],
) -> list[tuple[str, list[str]]]:
    """
    Filter (verse_text, stanza_tags) pairs by tradition tag.

    A stanza without its own [Tags: ...] marker inherits the file-level tags.
    With --include, a verse must share at least one tag with the include set.
    With --exclude, a verse is dropped if it shares any tag with the exclude set.
    When neither list is provided all verses are returned unchanged.
    """
    if not include and not exclude:
        return verses
    inc = {t.lower() for t in include}
    exc = {t.lower() for t in exclude}
    result = []
    for verse_text, stanza_tags in verses:
        tags = {t.lower() for t in (stanza_tags if stanza_tags else file_tags)}
        if inc and not tags & inc:
            continue
        if exc and tags & exc:
            continue
        result.append((verse_text, stanza_tags))
    return result


def _word_wrap(text: str, max_chars: int) -> list[str]:
    """Word-wrap text into lines of at most max_chars characters."""
    words = text.split()
    if not words:
        return ['']
    lines: list[str] = []
    current = ''
    for word in words:
        if not current:
            current = word
        elif len(current) + 1 + len(word) <= max_chars:
            current += ' ' + word
        else:
            lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines or ['']


def verse_to_pairs(verse: str, mel_lines: list[str],
                   max_lyric_chars: int = 25) -> list[tuple[str, str]]:
    """
    Divide verse text across the melody lines then word-wrap each portion
    to max_lyric_chars.  Returns a flat list of (melody_segment, lyric_line)
    pairs where every lyric line is ≤ max_lyric_chars characters.

    The melody segment for each phrase is repeated above every wrapped
    sub-line that belongs to it, so the "appropriate melody" always appears
    directly above each line of lyrics.
    """
    if not mel_lines:
        return [('', line) for line in _word_wrap(verse, max_lyric_chars)]

    words = verse.split()
    n     = len(mel_lines)
    per   = len(words) / n

    pairs: list[tuple[str, str]] = []
    for i, melody in enumerate(mel_lines):
        start = round(i * per)
        end   = min(round((i + 1) * per), len(words))
        group = ' '.join(words[start:end])
        for sub in _word_wrap(group, max_lyric_chars):
            pairs.append((melody, sub))

    return pairs


# ---------------------------------------------------------------------------
# Slide model
# ---------------------------------------------------------------------------

@dataclass
class Slide:
    kind:        str
    title:       str = ''
    attribution: str = ''
    verse_num:   int = 0
    pairs: list = field(default_factory=list)   # [(melody_str, lyric_str)]


def build_slides(hymn: dict, lines_per_slide: int = 3,
                 max_lyric_chars: int = 25,
                 include_tags: list[str] | None = None,
                 exclude_tags: list[str] | None = None) -> list[Slide]:
    slides: list[Slide] = [
        Slide('title', title=hymn['title'], attribution=hymn['attribution'])
    ]

    abc_lines  = _abc_body_lines(hymn['abc_text'])
    bar_counts = _barlines_per_line(abc_lines)
    mel_lines  = split_musiqwik(hymn['musiqwik'], bar_counts)

    if not mel_lines:
        print('[warn] No melody lines found; slides will contain lyrics only.',
              file=sys.stderr)

    all_verses = parse_verses(hymn['lyrics_text'])
    verses = filter_verses(
        all_verses,
        hymn.get('file_tags', []),
        include_tags or [],
        exclude_tags or [],
    )

    for v_idx, (verse, _stanza_tags) in enumerate(verses):
        pairs = verse_to_pairs(verse, mel_lines, max_lyric_chars)

        for start in range(0, max(len(pairs), 1), lines_per_slide):
            slides.append(Slide(
                'content',
                title=hymn['title'],
                verse_num=v_idx + 1,
                pairs=pairs[start : start + lines_per_slide],
            ))

    return slides


# ---------------------------------------------------------------------------
# HTML output
# ---------------------------------------------------------------------------

_CSS = """
* { box-sizing: border-box; margin: 0; padding: 0; }
html, body { background: #0d0d1a; }
.slide {
    width: 1280px; height: 720px;
    display: flex; flex-direction: column;
    justify-content: center; align-items: flex-start;
    padding: 56px 80px;
    background: #1a1a2e;
    page-break-after: always; page-break-inside: avoid;
    overflow: hidden;
}
.title-slide { align-items: center; text-align: center; gap: 28px; }
.title-slide h1 {
    font-family: 'OpenDyslexic Mono', monospace;
    font-size: 48px; font-weight: bold; color: #f0e68c; line-height: 1.25;
}
.title-slide .attr {
    font-family: 'OpenDyslexic Mono', monospace;
    font-size: 18px; color: #9090b8; white-space: pre-line;
    max-width: 960px; line-height: 1.65;
}
.verse-label {
    font-family: 'OpenDyslexic Mono', monospace;
    font-size: 15px; color: #6868a0; margin-bottom: 18px;
}
.pair { width: 100%; margin-bottom: 26px; }
.melody {
    font-family: 'Musiqwik', monospace;
    font-size: 38px; line-height: 1; color: #f0e68c;
    white-space: pre; letter-spacing: 0;
}
.lyric {
    font-family: 'OpenDyslexic Mono', monospace;
    font-size: 26px; line-height: 1.35; color: #e4e4e4;
    margin-top: 6px;
}
@media print {
    body { background: white; }
    .slide {
        background: white;
        page-break-after: always; page-break-inside: avoid;
        width: 100vw; height: 100vh;
    }
    .title-slide h1, .melody { color: #111; }
    .title-slide .attr, .verse-label { color: #555; }
    .lyric { color: #222; }
}
"""


# ---------------------------------------------------------------------------
# OTP (LibreOffice Impress) helpers
# ---------------------------------------------------------------------------

# Slide geometry in EMU (matches the PPTX output dimensions)
_W  = 9_144_000
_H  = 5_143_500
_MX = int(_W * 0.06)
_MY = int(_H * 0.08)


def _ecm(emu: int) -> str:
    """Convert EMU to a centimetre string for ODF attributes."""
    return f"{emu / 914400 * 2.54:.4f}cm"


_OTP_STYLES = """\
<?xml version="1.0" encoding="UTF-8"?>
<office:document-styles
  xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
  xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"
  xmlns:fo="urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0"
  xmlns:draw="urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
  xmlns:presentation="urn:oasis:names:tc:opendocument:xmlns:presentation:1.0"
  office:version="1.3">
  <office:styles>
    <style:style style:name="HymnPage" style:family="drawing-page">
      <style:drawing-page-properties draw:fill="solid" draw:fill-color="#1A1A2E"
        draw:background-size="border"
        presentation:display-header="false" presentation:display-footer="false"
        presentation:display-page-number="false" presentation:display-date-time="false"/>
    </style:style>
    <style:style style:name="P-Melody" style:family="paragraph">
      <style:paragraph-properties fo:text-align="left"/>
      <style:text-properties fo:font-family="Musiqwik" fo:font-size="26pt" fo:color="#F0E68C"/>
    </style:style>
    <style:style style:name="P-Lyrics" style:family="paragraph">
      <style:paragraph-properties fo:text-align="left"/>
      <style:text-properties fo:font-family="OpenDyslexic Mono" fo:font-size="19pt" fo:color="#E4E4E4"/>
    </style:style>
    <style:style style:name="P-Title" style:family="paragraph">
      <style:paragraph-properties fo:text-align="center"/>
      <style:text-properties fo:font-family="OpenDyslexic Mono" fo:font-size="36pt"
        fo:font-weight="bold" fo:color="#F0E68C"/>
    </style:style>
    <style:style style:name="P-Attr" style:family="paragraph">
      <style:paragraph-properties fo:text-align="center"/>
      <style:text-properties fo:font-family="OpenDyslexic Mono" fo:font-size="15pt" fo:color="#9090B8"/>
    </style:style>
    <style:style style:name="P-Label" style:family="paragraph">
      <style:paragraph-properties fo:text-align="left"/>
      <style:text-properties fo:font-family="OpenDyslexic Mono" fo:font-size="13pt" fo:color="#6868A0"/>
    </style:style>
  </office:styles>
  <office:master-styles>
    <style:master-page style:name="Default" draw:style-name="HymnPage"/>
  </office:master-styles>
</office:document-styles>"""


# ---------------------------------------------------------------------------
# HTML helpers
# ---------------------------------------------------------------------------

def _e(s: str) -> str:
    return (s.replace('&', '&amp;').replace('<', '&lt;')
             .replace('>', '&gt;').replace('"', '&quot;'))


def to_html(slides: list[Slide]) -> str:
    title = _e(slides[0].title) if slides else ''
    out   = [
        '<!DOCTYPE html>', '<html lang="en">', '<head>',
        '<meta charset="UTF-8">',
        '<meta name="viewport" content="width=1280">',
        f'<title>{title}</title>',
        f'<style>{_CSS}</style>',
        '</head>', '<body>',
    ]

    for sl in slides:
        if sl.kind == 'title':
            attr_html = _e(sl.attribution).replace('\n', '<br>')
            out += [
                '<div class="slide title-slide">',
                f'<h1>{_e(sl.title)}</h1>',
                f'<div class="attr">{attr_html}</div>',
                '</div>',
            ]
        else:
            out += [
                '<div class="slide">',
                f'<p class="verse-label">Verse {sl.verse_num}</p>',
            ]
            for melody, lyric in sl.pairs:
                out += [
                    '<div class="pair">',
                    f'<div class="melody">{_e(melody)}</div>',
                    f'<div class="lyric">{_e(lyric)}</div>',
                    '</div>',
                ]
            out.append('</div>')

    out += ['</body>', '</html>']
    return '\n'.join(out)


# ---------------------------------------------------------------------------
# PDF output (Playwright + pre-installed Chromium)
# ---------------------------------------------------------------------------

_CHROMIUM = '/opt/pw-browsers/chromium'


def to_pdf(html: str, out: Path) -> None:
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        sys.exit('PDF requires playwright: pip install playwright')

    with sync_playwright() as pw:
        browser = pw.chromium.launch(executable_path=_CHROMIUM)
        page    = browser.new_page()
        page.set_content(html, wait_until='domcontentloaded')
        page.pdf(path=str(out), width='1280px', height='720px',
                 print_background=True)
        browser.close()


# ---------------------------------------------------------------------------
# PPTX output (python-pptx)
# ---------------------------------------------------------------------------

def to_pptx(slides: list[Slide], out: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        sys.exit('PPTX requires python-pptx: pip install python-pptx')

    DARK   = RGBColor(0x1a, 0x1a, 0x2e)
    GOLD   = RGBColor(0xF0, 0xE6, 0x8C)
    WHITE  = RGBColor(0xE4, 0xE4, 0xE4)
    DIM    = RGBColor(0x90, 0x90, 0xB8)

    prs = Presentation()
    prs.slide_width  = Emu(9_144_000)  # 10 in
    prs.slide_height = Emu(5_143_500)  # 5.625 in  (16:9)

    W  = int(prs.slide_width)
    H  = int(prs.slide_height)
    MX = int(W * 0.06)
    MY = int(H * 0.08)

    blank = prs.slide_layouts[6]

    def _bg(sl):
        f = sl.background.fill
        f.solid()
        f.fore_color.rgb = DARK

    def _box(sl, left, top, width, height, lines_: list[tuple],
             fname: str, fsize: float, color: RGBColor,
             bold: bool = False, align=None, wrap: bool = True):
        """Add a text box; lines_ is a list of (text, font_name, font_size, color)
        or a single (text, …) tuple used for every paragraph."""
        from pptx.util import Pt
        box = sl.shapes.add_textbox(left, top, width, height)
        tf  = box.text_frame
        tf.word_wrap = wrap
        for i, (txt, fn, fs, clr) in enumerate(lines_):
            p  = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if align:
                p.alignment = align
            run = p.add_run()
            run.text           = txt
            run.font.name      = fn
            run.font.size      = Pt(fs)
            run.font.color.rgb = clr
            run.font.bold      = bold
        return box

    def _tb(sl, left, top, width, height, text: str,
            fname: str, fsize: float, color: RGBColor,
            bold=False, align=None, wrap=True):
        """Single-font, multi-paragraph text box (splits on newlines)."""
        text_lines = text.splitlines() or ['']
        rows = [(t, fname, fsize, color) for t in text_lines]
        return _box(sl, left, top, width, height, rows,
                    fname, fsize, color, bold=bold, align=align, wrap=wrap)

    for data in slides:
        sl = prs.slides.add_slide(blank)
        _bg(sl)

        if data.kind == 'title':
            from pptx.enum.text import PP_ALIGN
            _tb(sl, MX, int(H * 0.20), W - 2*MX, int(H * 0.28),
                data.title, 'OpenDyslexic Mono', 36, GOLD,
                bold=True, align=PP_ALIGN.CENTER)
            _tb(sl, MX, int(H * 0.52), W - 2*MX, int(H * 0.42),
                data.attribution, 'OpenDyslexic Mono', 15, DIM,
                align=PP_ALIGN.CENTER)
        else:
            from pptx.enum.text import PP_ALIGN
            _tb(sl, MX, int(MY * 0.4), W - 2*MX, int(H * 0.08),
                f'Verse {data.verse_num}', 'OpenDyslexic Mono', 13, DIM)

            n      = len(data.pairs)
            avail  = H - MY - int(H * 0.12)
            pair_h = avail // max(n, 1)
            mel_h  = int(pair_h * 0.44)
            lyr_h  = int(pair_h * 0.50)

            for i, (melody, lyric) in enumerate(data.pairs):
                y = MY + int(H * 0.10) + i * pair_h
                _tb(sl, MX, y,          W - 2*MX, mel_h,
                    melody, 'Musiqwik', 26, GOLD, wrap=False)
                _tb(sl, MX, y + mel_h,  W - 2*MX, lyr_h,
                    lyric,  'OpenDyslexic Mono', 19, WHITE)

    prs.save(str(out))


# ---------------------------------------------------------------------------
# OTP output (LibreOffice Impress template)
# ---------------------------------------------------------------------------

def to_otp(slides: list[Slide], out: Path) -> None:
    """Write a LibreOffice Impress template (.otp) file with all hymn slides."""

    def _frame(x: int, y: int, w: int, h: int, pstyle: str,
                lines: list[str]) -> str:
        paras = ''.join(
            f'<text:p text:style-name="{pstyle}">{_xe(l)}</text:p>'
            for l in lines
        )
        return (
            f'<draw:frame draw:style-name="GrText" draw:layer="layout"'
            f' svg:x="{_ecm(x)}" svg:y="{_ecm(y)}"'
            f' svg:width="{_ecm(w)}" svg:height="{_ecm(h)}">'
            f'<draw:text-box>{paras}</draw:text-box></draw:frame>'
        )

    def _title_page(sl: Slide) -> str:
        fw = _W - 2 * _MX
        return (
            '<draw:page draw:name="Title" draw:style-name="dp1"'
            ' draw:master-page-name="Default">'
            + _frame(_MX, int(_H * 0.20), fw, int(_H * 0.28), 'P-Title', [sl.title])
            + _frame(_MX, int(_H * 0.52), fw, int(_H * 0.42), 'P-Attr',
                     sl.attribution.splitlines() or [''])
            + '</draw:page>'
        )

    def _content_page(sl: Slide, idx: int) -> str:
        n      = max(len(sl.pairs), 1)
        avail  = _H - _MY - int(_H * 0.12)
        pair_h = avail // n
        mel_h  = int(pair_h * 0.44)
        lyr_h  = int(pair_h * 0.50)
        fw     = _W - 2 * _MX
        page   = (
            f'<draw:page draw:name="{_xe(f"V{sl.verse_num}-{idx}")}"'
            ' draw:style-name="dp1" draw:master-page-name="Default">'
        )
        page += _frame(_MX, int(_MY * 0.4), fw, int(_H * 0.08),
                       'P-Label', [f'Verse {sl.verse_num}'])
        for i, (melody, lyric) in enumerate(sl.pairs):
            y = _MY + int(_H * 0.10) + i * pair_h
            page += _frame(_MX, y,         fw, mel_h, 'P-Melody', [melody])
            page += _frame(_MX, y + mel_h, fw, lyr_h, 'P-Lyrics', [lyric])
        page += '</draw:page>'
        return page

    pages: list[str] = []
    verse_counters: dict[int, int] = {}
    for sl in slides:
        if sl.kind == 'title':
            pages.append(_title_page(sl))
        else:
            verse_counters[sl.verse_num] = verse_counters.get(sl.verse_num, 0) + 1
            pages.append(_content_page(sl, verse_counters[sl.verse_num]))

    content_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<office:document-content'
        ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
        ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
        ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0"'
        ' xmlns:draw="urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"'
        ' xmlns:svg="urn:oasis:names:tc:opendocument:xmlns:svg-compatible:1.0"'
        ' xmlns:fo="urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0"'
        ' xmlns:presentation="urn:oasis:names:tc:opendocument:xmlns:presentation:1.0"'
        ' office:version="1.3">'
        '<office:automatic-styles>'
        '<style:style style:name="dp1" style:family="drawing-page">'
        '<style:drawing-page-properties draw:fill="solid" draw:fill-color="#1A1A2E"'
        ' draw:background-size="border"'
        ' presentation:display-header="false" presentation:display-footer="false"'
        ' presentation:display-page-number="false"'
        ' presentation:display-date-time="false"/>'
        '</style:style>'
        '<style:style style:name="GrText" style:family="graphic">'
        '<style:graphic-properties draw:fill="none" draw:stroke="none"/>'
        '</style:style>'
        '</office:automatic-styles>'
        '<office:body><office:presentation>'
        + ''.join(pages)
        + '</office:presentation></office:body></office:document-content>'
    )

    manifest = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<manifest:manifest'
        ' xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0"'
        ' manifest:version="1.3">'
        '<manifest:file-entry manifest:full-path="/"'
        ' manifest:media-type="application/vnd.oasis.opendocument.presentation-template"'
        ' manifest:version="1.3"/>'
        '<manifest:file-entry manifest:full-path="content.xml"'
        ' manifest:media-type="text/xml"/>'
        '<manifest:file-entry manifest:full-path="styles.xml"'
        ' manifest:media-type="text/xml"/>'
        '</manifest:manifest>'
    )

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(
            zipfile.ZipInfo('mimetype'),
            b'application/vnd.oasis.opendocument.presentation-template',
            compress_type=zipfile.ZIP_STORED,
        )
        zf.writestr('META-INF/manifest.xml', manifest)
        zf.writestr('styles.xml', _OTP_STYLES)
        zf.writestr('content.xml', content_xml)
    out.write_bytes(buf.getvalue())


# ---------------------------------------------------------------------------
# ProPresenter 6 output (.pro6)
# ---------------------------------------------------------------------------

def to_pro6(slides: list[Slide], out: Path) -> None:
    """Write a ProPresenter 6 (.pro6) XML file with all hymn slides."""

    PW, PH = 1280, 720   # slide pixel dimensions

    def _uid() -> str:
        return str(uuid.uuid4())

    def _rtf(text: str, font: str, pt: int, r: int, g: int, b: int) -> str:
        fs = pt * 2   # RTF uses half-points
        escaped = text.replace('\\', '\\\\').replace('{', '\\{').replace('}', '\\}')
        lines   = escaped.splitlines() or ['']
        body    = '\\par\n'.join(lines)
        return (
            f'{{\\rtf1\\ansi\\deff0'
            f'{{\\fonttbl{{\\f0\\fnil {font};}}}}'
            f'{{\\colortbl ;\\red{r}\\green{g}\\blue{b};}}'
            f'\\pard\\f0\\fs{fs}\\cf1 {body}}}'
        )

    def _b64rtf(text: str, font: str, pt: int, r: int, g: int, b: int) -> str:
        return base64.b64encode(_rtf(text, font, pt, r, g, b).encode('latin-1')).decode()

    def _tel(name: str, x: int, y: int, w: int, h: int, rtf64: str) -> str:
        return (
            f'<RVTextElement additionalLineFillHeight="0" adjustsHeightToFit="false"'
            f' bezelRadius="0" displayDelay="0" displayName="{_xe(name)}"'
            f' drawLineBackground="false" drawingFill="false" drawingShadow="false"'
            f' drawingStroke="false" fillColor="0 0 0 0" fromTemplate="false"'
            f' guid="{_uid()}" locked="false" opacity="1" persistent="0"'
            f' revealType="0" rotation="0" scaleFactor="0"'
            f' shadow="0 0 0 0 0" stroke="0 0 0 0 1" textSourceRemote="0"'
            f' typeID="0" useAllCaps="false">'
            f'<_-RVRect3D-_position x="{float(x)}" y="{float(y)}" z="0.0"/>'
            f'<_-RVRect3D-_size x="{float(w)}" y="{float(h)}" z="0.0"/>'
            f'<RTFData>{rtf64}</RTFData>'
            f'</RVTextElement>'
        )

    def _dslide(elements: list[str]) -> str:
        return (
            f'<RVDisplaySlide backgroundColor="0 0 0 1" enabled="true"'
            f' highlightColor="0 0 0 0" hotKey="" label="" notes=""'
            f' uuid="{_uid()}"><cues/>'
            f'<displayElements>{"".join(elements)}</displayElements>'
            f'</RVDisplaySlide>'
        )

    def _group(name: str, dslides: list[str]) -> str:
        return (
            f'<RVSlideGrouping color="0 0 0 0" name="{_xe(name)}" uuid="{_uid()}">'
            f'<slides>{"".join(dslides)}</slides></RVSlideGrouping>'
        )

    title_sl  = slides[0] if slides and slides[0].kind == 'title' else None
    song_title = title_sl.title if title_sl else ''
    attr       = title_sl.attribution if title_sl else ''

    t_el = _tel('Title', 50, int(PH*0.20), PW-100, int(PH*0.28),
                _b64rtf(song_title, 'OpenDyslexic Mono', 36, 0xF0, 0xE6, 0x8C))
    a_el = _tel('Attribution', 50, int(PH*0.52), PW-100, int(PH*0.38),
                _b64rtf(attr, 'OpenDyslexic Mono', 15, 0x90, 0x90, 0xB8))
    groups = [_group('Title', [_dslide([t_el, a_el])])]

    verse_slides: dict[int, list[Slide]] = defaultdict(list)
    for sl in slides:
        if sl.kind == 'content':
            verse_slides[sl.verse_num].append(sl)

    for v_num in sorted(verse_slides.keys()):
        dslides = []
        for sl in verse_slides[v_num]:
            n      = max(len(sl.pairs), 1)
            avail  = PH - 80 - 40   # top area for label + bottom margin
            pair_h = avail // n
            mel_h  = max(int(pair_h * 0.44), 1)
            lyr_h  = max(int(pair_h * 0.50), 1)
            els = [_tel('Label', 50, 40, PW-100, 25,
                        _b64rtf(f'Verse {v_num}', 'OpenDyslexic Mono', 13,
                                0x68, 0x68, 0xA0))]
            for i, (melody, lyric) in enumerate(sl.pairs):
                y = 80 + i * pair_h
                els.append(_tel('Melody', 50, y, PW-100, mel_h,
                                _b64rtf(melody, 'Musiqwik', 26, 0xF0, 0xE6, 0x8C)))
                els.append(_tel('Lyric', 50, y+mel_h, PW-100, lyr_h,
                                _b64rtf(lyric, 'OpenDyslexic Mono', 19, 0xE4, 0xE4, 0xE4)))
            dslides.append(_dslide(els))
        groups.append(_group(f'Verse {v_num}', dslides))

    out.write_text(
        '<?xml version="1.0" encoding="utf-8"?>\n'
        f'<RVPresentationDocument CCLISongTitle="{_xe(song_title)}"'
        f' CCLICopyrightInfo="" CCLIArtistCredits=""'
        f' backgroundColor="0 0 0 1" buildNumber="6105" category="Presentation"'
        f' docType="0" drawingBackgroundColor="false" height="{PH}"'
        f' lastDateUsed="2026-06-28T00:00:00" notes="" resourcesDirectory=""'
        f' selectedArrangement="" usedCount="0" uuid="{_uid()}"'
        f' versionNumber="600" width="{PW}">'
        f'<groups>{"".join(groups)}</groups>'
        f'</RVPresentationDocument>\n',
        encoding='utf-8',
    )


# ---------------------------------------------------------------------------
# EasyWorship Song XML output (.ewsx)
# ---------------------------------------------------------------------------

def to_ewsx(slides: list[Slide], out: Path) -> None:
    """Write an EasyWorship Song XML (.ewsx) file.

    Lyrics are organised by verse.  Melody characters appear as the first
    line of each verse so users can manually apply the Musiqwik font in
    EasyWorship's song editor.  For slides with melody rendered
    automatically use --format pptx and import into EasyWorship.
    """
    title_sl = slides[0] if slides and slides[0].kind == 'title' else None
    title    = title_sl.title if title_sl else ''
    attr     = title_sl.attribution if title_sl else ''

    author = ''
    copyright_str = 'Public Domain'
    for line in attr.splitlines():
        ll = line.lower()
        if ll.startswith('words:'):
            author = line[6:].strip().rstrip('.')
        elif ll.startswith('copyright:'):
            copyright_str = line[10:].strip().rstrip('.')

    verse_slides: dict[int, list[Slide]] = defaultdict(list)
    for sl in slides:
        if sl.kind == 'content':
            verse_slides[sl.verse_num].append(sl)

    verse_ids    = sorted(verse_slides.keys())
    verse_order  = ' '.join(f'V{v}' for v in verse_ids)

    def _verse_lines(v: int) -> str:
        parts: list[str] = []
        for sl in verse_slides[v]:
            for melody, lyric in sl.pairs:
                if melody:
                    parts.append(_xe(melody))
                if lyric:
                    parts.append(_xe(lyric))
        return '&#10;'.join(parts)

    verses_xml = '\n'.join(
        f'        <verse id="V{v}" label="Verse {v}">'
        f'<lines>{_verse_lines(v)}</lines></verse>'
        for v in verse_ids
    )

    notes = (
        'Melody lines use Musiqwik font characters. '
        'In EasyWorship Song Editor select each melody line and apply the Musiqwik font. '
        'For slides with melody rendered automatically use '
        'hymn_to_presentation.py --format pptx and import into EasyWorship.'
    )

    out.write_text(
        '<?xml version="1.0" encoding="utf-8"?>\n'
        '<EasywrshipSongData version="3">\n'
        '  <songs>\n'
        '    <song>\n'
        f'      <title>{_xe(title)}</title>\n'
        f'      <author>{_xe(author)}</author>\n'
        f'      <copyright>{_xe(copyright_str)}</copyright>\n'
        '      <ccli/>\n'
        f'      <notes>{_xe(notes)}</notes>\n'
        f'      <verseOrder>{verse_order}</verseOrder>\n'
        '      <lyrics>\n'
        f'{verses_xml}\n'
        '      </lyrics>\n'
        '    </song>\n'
        '  </songs>\n'
        '</EasywrshipSongData>\n',
        encoding='utf-8',
    )


# ---------------------------------------------------------------------------
# Tradition-filter prompt
# ---------------------------------------------------------------------------

def _prompt_tradition_filter(
    hymn_title: str,
    file_tags: list[str],
    verses: list[tuple[str, list[str]]],
) -> tuple[list[str], list[str]]:
    """
    When the hymn contains stanzas with tradition-specific tags, describe them
    and interactively ask the user how to filter.

    A stanza is "notable" when its explicit [Tags: ...] marker differs from
    the file-level tag set — meaning it was added by, or is exclusive to, a
    particular tradition within the broader scope of the hymn file.

    Returns (include_tags, exclude_tags).  Both empty means include all.
    When stdin is not a TTY, prints a one-line notice to stderr and returns
    ([], []) so piped / scripted use is unaffected.
    """
    file_set = frozenset(t.lower() for t in file_tags)

    notable: list[tuple[int, list[str], str]] = []
    for i, (text, tags) in enumerate(verses):
        if not tags:
            continue
        if frozenset(t.lower() for t in tags) != file_set:
            preview = (text[:60] + '…') if len(text) > 60 else text
            notable.append((i + 1, tags, preview))

    if not notable:
        return [], []

    if not sys.stdin.isatty():
        print(
            f'[info] {hymn_title}: {len(notable)} stanza(s) carry tradition-specific '
            'tags; use --include / --exclude to filter. Proceeding with all stanzas.',
            file=sys.stderr,
        )
        return [], []

    # Collect unique tradition sets from notable stanzas (order of first occurrence).
    seen_sets: list[frozenset] = []
    for _, tags, _ in notable:
        s = frozenset(t.lower() for t in tags)
        if s not in seen_sets:
            seen_sets.append(s)

    # Build a numbered menu — one "include only" entry per unique tradition set.
    menu: list[tuple[str, str, list[str]]] = []   # (key, label, include_tags)
    for idx, tset in enumerate(seen_sets, start=1):
        label_tags = ', '.join(sorted(tset))
        count = sum(
            1 for _, stanza_tags in verses
            if (frozenset(t.lower() for t in stanza_tags) if stanza_tags else file_set) & tset
        )
        noun = 'stanza' if count == 1 else 'stanzas'
        menu.append((str(idx), f'{label_tags.capitalize()} only  ({count} {noun})', sorted(tset)))

    # Print header and notable-stanza summary.
    bar = '─' * 62
    print(f'\n{bar}')
    print(f'  {hymn_title}')
    print(f'  File tradition: {", ".join(file_tags) or "(none)"}')
    print(bar)
    print('\n  Stanzas with tradition-specific tags:\n')
    for stanza_num, tags, preview in notable:
        print(f'    Stanza {stanza_num}  [{", ".join(tags)}]')
        print(f'    "{preview}"\n')

    total = len(verses)
    noun = 'stanza' if total == 1 else 'stanzas'
    print('  Options:\n')
    print(f'    a  Include all  ({total} {noun}, default)')
    for key, label, _ in menu:
        print(f'    {key}  {label}')
    print('    c  Custom filter')
    print()

    while True:
        try:
            choice = input('  Choice [a]: ').strip().lower() or 'a'
        except (EOFError, KeyboardInterrupt):
            print()
            return [], []

        if choice == 'a':
            return [], []

        for key, _, inc in menu:
            if choice == key:
                return inc, []

        if choice == 'c':
            inc_raw = input('  Include tags (comma-separated, blank = all): ').strip()
            exc_raw = input('  Exclude tags (comma-separated, blank = none): ').strip()
            inc = [t.strip() for t in inc_raw.split(',') if t.strip()]
            exc = [t.strip() for t in exc_raw.split(',') if t.strip()]
            return inc, exc

        print('  Please enter one of the listed choices.  ', end='', flush=True)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main() -> None:
    ap = argparse.ArgumentParser(
        description='Convert a hymn file to HTML, PDF, or PPTX presentation slides.',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    ap.add_argument('--file', '-f', required=True,
                    help='Path to the hymn file')
    ap.add_argument('--format', choices=['html', 'pdf', 'pptx', 'otp', 'pro6', 'ewsx'],
                    default='html', help='Output format (default: html)')
    ap.add_argument('--output', '-o',
                    help='Output file path (default: <hymn_name>.<format>)')
    ap.add_argument('--lines-per-slide', type=int, default=3, metavar='N',
                    help='Max melody/lyric pairs per content slide (default: 3)')
    ap.add_argument('--max-lyric-chars', type=int, default=25, metavar='C',
                    help='Max characters per lyric line (default: 25)')
    ap.add_argument('--include', '-I', action='append', default=[], metavar='TAG',
                    help='Include only stanzas with this tradition tag (repeatable)')
    ap.add_argument('--exclude', '-X', action='append', default=[], metavar='TAG',
                    help='Exclude stanzas with this tradition tag (repeatable)')
    args = ap.parse_args()

    src = Path(args.file)
    if not src.exists():
        sys.exit(f'Error: file not found: {src}')

    hymn = parse_hymn_file(src)

    if not args.include and not args.exclude:
        args.include, args.exclude = _prompt_tradition_filter(
            hymn['title'],
            hymn.get('file_tags', []),
            parse_verses(hymn['lyrics_text']),
        )

    slides = build_slides(hymn, args.lines_per_slide, args.max_lyric_chars,
                          args.include, args.exclude)
    out    = Path(args.output) if args.output else Path(src.name + '.' + args.format)

    if args.format == 'html':
        html = to_html(slides)
        out.write_text(html, encoding='utf-8')
    elif args.format == 'pdf':
        to_pdf(to_html(slides), out)
    elif args.format == 'pptx':
        to_pptx(slides, out)
    elif args.format == 'otp':
        to_otp(slides, out)
    elif args.format == 'pro6':
        to_pro6(slides, out)
    elif args.format == 'ewsx':
        to_ewsx(slides, out)

    print(f'Written {len(slides)} slides → {out}')


if __name__ == '__main__':
    main()
